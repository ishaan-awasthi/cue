"""RAG pipeline.

Two responsibilities:
1. File ingestion: extract text from uploaded files, chunk it (~500 tokens),
   embed with OpenAI text-embedding-3-small, store in Supabase pgvector.
2. Q&A retrieval: embed an incoming question, similarity-search pgvector,
   pass question + top chunks to chat model (Anthropic > OpenRouter > OpenAI),
   return RAGResult with answer, confidence, sources, fallback flag.
"""

from __future__ import annotations

import asyncio
import io
import logging
from dataclasses import dataclass
from typing import Any

import tiktoken
from openai import AsyncOpenAI

from ..config import settings
from ..db import queries
from ..db.models import DocumentChunk

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Types
# ---------------------------------------------------------------------------

@dataclass
class RAGResult:
    answer: str
    confidence: float
    sources_used: list[str]
    fallback_used: bool
    source: str
    supporting_context: list[dict[str, str]]


# ---------------------------------------------------------------------------
# Clients (lazily shared)
# ---------------------------------------------------------------------------

_openai: AsyncOpenAI | None = None
_anthropic_client = None
_chat_client = None

EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIM = 1536

CHUNK_TOKENS = 500
CHUNK_OVERLAP = 50  # token overlap between adjacent chunks

_TOKENIZER = None


def _get_tokenizer():
    global _TOKENIZER
    if _TOKENIZER is not None:
        return _TOKENIZER
    try:
        _TOKENIZER = tiktoken.get_encoding("cl100k_base")
    except Exception:
        _TOKENIZER = False
    return _TOKENIZER


def _get_openai() -> AsyncOpenAI:
    global _openai
    if _openai is None:
        _openai = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
    return _openai


def _get_anthropic():
    """Lazy init Anthropic client when ANTHROPIC_API_KEY is set."""
    global _anthropic_client
    if _anthropic_client is None and settings.ANTHROPIC_API_KEY:
        from anthropic import AsyncAnthropic
        _anthropic_client = AsyncAnthropic(api_key=settings.ANTHROPIC_API_KEY)
    return _anthropic_client


def _get_chat_client():
    """OpenRouter or OpenAI for chat when Anthropic not used."""
    global _chat_client
    if _chat_client is None:
        if settings.OPENROUTER_API_KEY:
            _chat_client = ("openrouter", AsyncOpenAI(
                base_url="https://openrouter.ai/api/v1",
                api_key=settings.OPENROUTER_API_KEY,
            ))
        else:
            _chat_client = ("openai", _get_openai())
    return _chat_client


# ---------------------------------------------------------------------------
# Chat completion — single entry point; priority: Anthropic > OpenRouter > OpenAI
# ---------------------------------------------------------------------------

async def _chat_complete(prompt: str, max_tokens: int = 150) -> str:
    """Generate completion. Uses Anthropic if set, else OpenRouter, else OpenAI."""
    anthropic = _get_anthropic()
    if anthropic:
        try:
            msg = await anthropic.messages.create(
                model="claude-3-5-haiku-20241022",
                max_tokens=max_tokens,
                messages=[{"role": "user", "content": prompt}],
            )
            if msg.content and len(msg.content) > 0:
                block = msg.content[0]
                if hasattr(block, "text"):
                    return block.text.strip()
            return ""
        except Exception as exc:
            logger.warning("Anthropic chat failed, falling back: %s", exc)

    provider, client = _get_chat_client()
    model = "openai/gpt-4o-mini" if provider == "openrouter" else "gpt-4o-mini"
    resp = await client.chat.completions.create(
        model=model,
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    )
    return (resp.choices[0].message.content or "").strip()


# ---------------------------------------------------------------------------
# Embedding helpers (OpenAI only)
# ---------------------------------------------------------------------------

async def embed_text(text: str) -> list[float]:
    """Return a 1536-dim embedding for *text* using OpenAI."""
    resp = await _get_openai().embeddings.create(model=EMBEDDING_MODEL, input=text)
    return resp.data[0].embedding


# ---------------------------------------------------------------------------
# File ingestion
# ---------------------------------------------------------------------------

async def ingest_file(
    file_id: str,
    user_id: str,
    filename: str,
    file_type: str,
    content: bytes,
) -> int:
    """Extract, chunk, embed, and store the file. Returns chunk count. Returns 0 on embedding failure (no 500)."""
    try:
        sections = _extract_sections(filename, file_type, content)
        if not sections:
            return 0

        chunks = _chunk_sections(sections)
        stored = 0
        for idx, chunk in enumerate(chunks):
            try:
                embedding = await embed_text(chunk["text"])
            except Exception as exc:
                logger.warning("Embedding failed for chunk %d of %s: %s", idx + 1, filename, exc)
                continue
            queries.insert_chunk(
                file_id=file_id,
                user_id=user_id,
                chunk_text=chunk["text"],
                chunk_index=idx,
                embedding_vector=embedding,
                metadata=chunk["metadata"],
            )
            stored += 1

        logger.info("Ingested %d chunks for file %s", stored, filename)
        return stored
    except Exception as exc:
        logger.error("File ingestion failed for %s: %s", filename, exc)
        return 0


def _extract_sections(filename: str, file_type: str, content: bytes) -> list[dict[str, Any]]:
    """Return extraction units with best-effort location metadata."""
    ft = file_type.lower().lstrip(".")
    name_lower = filename.lower()

    if ft in ("pdf",) or name_lower.endswith(".pdf"):
        return _extract_pdf_sections(content)
    if ft in ("pptx",) or name_lower.endswith(".pptx"):
        return _extract_pptx_sections(content)
    if ft in ("docx",) or name_lower.endswith(".docx"):
        return _extract_docx_sections(content)
    if ft in ("md", "markdown") or name_lower.endswith(".md"):
        text = content.decode("utf-8", errors="replace").strip()
        return [{"text": text, "metadata": {"section": "markdown"}}] if text else []
    text = content.decode("utf-8", errors="replace").strip()
    return [{"text": text, "metadata": {"section": "text"}}] if text else []


def _extract_pdf_sections(content: bytes) -> list[dict[str, Any]]:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    sections: list[dict[str, Any]] = []
    for idx, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            sections.append({
                "text": text.strip(),
                "metadata": {"page_number": idx + 1, "location": f"Page {idx + 1}"},
            })
    return sections


def _extract_pptx_sections(content: bytes) -> list[dict[str, Any]]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    sections: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = " ".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        merged = "\n".join(parts).strip()
        if merged:
            sections.append({
                "text": merged,
                "metadata": {"slide_number": slide_idx + 1, "location": f"Slide {slide_idx + 1}"},
            })
    return sections


def _extract_docx_sections(content: bytes) -> list[dict[str, Any]]:
    from docx import Document
    doc = Document(io.BytesIO(content))
    sections: list[dict[str, Any]] = []
    current_section = "Document"
    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        if p.style and p.style.name and "heading" in p.style.name.lower():
            current_section = text
            continue
        sections.append({
            "text": text,
            "metadata": {"section": current_section, "location": f"Section: {current_section}"},
        })
    return sections


def _chunk_text_with_overlap(text: str) -> list[str]:
    tokenizer = _get_tokenizer()
    if tokenizer is False:
        # Offline fallback when tokenizer assets are unavailable.
        words = text.split()
        if not words:
            return []
        max_words = 350
        overlap = 40
        chunks: list[str] = []
        start = 0
        while start < len(words):
            end = min(start + max_words, len(words))
            chunks.append(" ".join(words[start:end]))
            if end == len(words):
                break
            start += max_words - overlap
        return chunks

    tokens = tokenizer.encode(text)
    chunks: list[str] = []
    start = 0
    while start < len(tokens):
        end = min(start + CHUNK_TOKENS, len(tokens))
        chunk_tokens = tokens[start:end]
        chunks.append(tokenizer.decode(chunk_tokens))
        if end == len(tokens):
            break
        start += CHUNK_TOKENS - CHUNK_OVERLAP
    return chunks


def _chunk_sections(sections: list[dict[str, Any]]) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    for section in sections:
        text = str(section.get("text", "")).strip()
        metadata = dict(section.get("metadata") or {})
        if not text:
            continue
        for part_idx, chunk_text in enumerate(_chunk_text_with_overlap(text)):
            merged_metadata = dict(metadata)
            merged_metadata["chunk_part"] = part_idx
            chunks.append({"text": chunk_text, "metadata": merged_metadata})
    return chunks


# ---------------------------------------------------------------------------
# RAG helpers
# ---------------------------------------------------------------------------

def _deduplicated_sources(chunks: list[DocumentChunk]) -> list[str]:
    """Unique ordered filenames from chunk list."""
    seen: set[str] = set()
    result: list[str] = []
    for c in chunks:
        fn = (c.filename or "").strip()
        if fn and fn not in seen:
            seen.add(fn)
            result.append(fn)
    return result


async def _retrieve_chunks(user_id: str, embedding: list[float], top_k: int = 5) -> list[DocumentChunk]:
    return await asyncio.to_thread(queries.similarity_search, user_id, embedding, top_k)


# ---------------------------------------------------------------------------
# Q&A retrieval
# ---------------------------------------------------------------------------

async def _gpt_fallback(question: str) -> RAGResult:
    """Answer without RAG context when chunks are missing or below threshold."""
    prompt = (
        f"You are a real-time presentation coach whispering a brief answer into the "
        f"speaker's earpiece. Answer this audience question in 1-3 concise sentences. "
        f"Do not use filler phrases — speak naturally.\n\n"
        f"Audience question: {question}\n\nAnswer:"
    )
    try:
        answer = await _chat_complete(prompt, max_tokens=150)
        return RAGResult(
            answer=answer,
            confidence=0.0,
            sources_used=[],
            fallback_used=True,
            source="llm_fallback",
            supporting_context=[],
        )
    except Exception as exc:
        logger.error("GPT fallback failed: %s", exc)
        return RAGResult(
            answer="",
            confidence=0.0,
            sources_used=[],
            fallback_used=True,
            source="llm_fallback",
            supporting_context=[],
        )


def _build_supporting_context(chunks: list[DocumentChunk]) -> list[dict[str, str]]:
    seen: set[tuple[str, str]] = set()
    out: list[dict[str, str]] = []
    for chunk in chunks:
        filename = (chunk.filename or "").strip()
        md = chunk.metadata or {}
        location = str(md.get("location") or md.get("section") or f"Chunk {chunk.chunk_index}")
        key = (filename, location)
        if key in seen:
            continue
        seen.add(key)
        out.append({"fileName": filename or "uploaded file", "location": location})
        if len(out) >= 3:
            break
    return out


async def answer_question(user_id: str, question: str, session_id: str | None = None) -> RAGResult:
    """Embed question, retrieve chunks, filter by QA_MIN_CHUNK_SIMILARITY, generate answer.
    Falls back to GPT-only when no chunks pass threshold or embedding fails."""
    try:
        query_embedding = await embed_text(question)
    except Exception as exc:
        logger.error("Embedding failed in answer_question: %s", exc)
        return await _gpt_fallback(question)

    chunks = await asyncio.to_thread(queries.similarity_search, user_id, query_embedding, 5, session_id)
    min_sim = settings.QA_MIN_CHUNK_SIMILARITY
    usable = [c for c in chunks if (c.similarity or 0) >= min_sim]

    if not usable:
        return await _gpt_fallback(question)

    confidence = usable[0].similarity or 0.0
    sources = _deduplicated_sources(usable)
    supporting_context = _build_supporting_context(usable)
    context = "\n\n---\n\n".join(c.chunk_text for c in usable)
    prompt = (
        f"You are a real-time presentation coach whispering a brief answer into the "
        f"speaker's earpiece. Based only on the reference material below, answer the "
        f"audience question in 1-2 concise sentences for a live VC meeting. Keep it direct, "
        f"confident, and practical. Do not say 'based on the material' or use filler phrases.\n\n"
        f"Reference material:\n{context}\n\n"
        f"Audience question: {question}\n\nAnswer:"
    )

    try:
        answer = await _chat_complete(prompt, max_tokens=150)
        return RAGResult(
            answer=answer.strip(),
            confidence=confidence,
            sources_used=sources,
            fallback_used=False,
            source="session_docs",
            supporting_context=supporting_context,
        )
    except Exception as exc:
        logger.error("RAG chat failed: %s", exc)
        return await _gpt_fallback(question)
